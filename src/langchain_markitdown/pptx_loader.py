from typing import List, Dict, Any, Optional
from langchain_core.documents import Document
from .base_loader import BaseMarkitdownLoader
from langchain_text_splitters import MarkdownHeaderTextSplitter
from langchain_core.language_models import BaseChatModel
import re
import os
import io
import logging
from .utils import langchain_caption_adapter  # Import the adapter function


class PptxLoader(BaseMarkitdownLoader):
    def __init__(self, file_path: str, split_by_page: bool = False, llm: Optional[BaseChatModel] = None, prompt: Optional[str] = None, verbose: Optional[bool] = None):  # verbose is optional
        super().__init__(file_path, verbose=verbose) # Use inherited or passed-in value
        self.split_by_page = split_by_page
        self.llm = llm
        self.prompt = prompt
        self.logger.info(f"Initialized PptxLoader for {file_path} with split_by_page={split_by_page}")  # Use instance logger
        self.logger.info(f"Langchain LLM for image captioning: {llm.__class__.__name__ if llm else 'None'}")  # Use instance logger

    def load(
        self, 
        headers_to_split_on: Optional[List[str]] = None
    ) -> List[Document]:
        """Load a PPTX file and convert it to Langchain documents, splitting by Markdown headers."""
        try:
            # Basic converter for fallback when MarkitdownConverterOptions isn't available
            from markitdown import MarkItDown
            
            self.logger.info(f"Starting to load PPTX file: {self.file_path}")  # Use instance logger
            
            # Create basic metadata
            metadata: Dict[str, Any] = {
                "source": self.file_path,
                "file_name": self._get_file_name(self.file_path),
                "file_size": self._get_file_size(self.file_path),
                "conversion_success": True,
            }
            
            # Extract detailed metadata from PPTX using python-pptx
            try:
                from pptx import Presentation
                prs = Presentation(self.file_path)
                
                self.logger.info(f"Extracting metadata from PPTX file")  # Use instance logger
                
                # Basic presentation stats
                metadata["slide_count"] = len(prs.slides)
                self.logger.info(f"Found {metadata['slide_count']} slides in the presentation")  # Use instance logger

                # Core properties
                if hasattr(prs, 'core_properties'):
                    core_props = prs.core_properties
                    if hasattr(core_props, 'author') and core_props.author:
                        metadata["author"] = core_props.author
                    if hasattr(core_props, 'title') and core_props.title:
                        metadata["title"] = core_props.title
                    if hasattr(core_props, 'subject') and core_props.subject:
                        metadata["subject"] = core_props.subject
                    if hasattr(core_props, 'keywords') and core_props.keywords:
                        metadata["keywords"] = core_props.keywords
                    if hasattr(core_props, 'created') and core_props.created:
                        metadata["created"] = str(core_props.created)
                    if hasattr(core_props, 'modified') and core_props.modified:
                        metadata["modified"] = str(core_props.modified)
                    if hasattr(core_props, 'last_modified_by') and core_props.last_modified_by:
                        metadata["last_modified_by"] = core_props.last_modified_by
                    if hasattr(core_props, 'category') and core_props.category:
                        metadata["category"] = core_props.category
                    if hasattr(core_props, 'revision') and core_props.revision:
                        metadata["revision"] = core_props.revision
                
                # Count media elements
                image_count = 0
                text_box_count = 0
                chart_count = 0
                table_count = 0
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'shape_type'):
                            if shape.shape_type == 13:  # Picture
                                image_count += 1
                            elif shape.shape_type == 17:  # TextBox
                                text_box_count += 1
                            elif shape.shape_type == 3:  # Chart
                                chart_count += 1
                            elif shape.shape_type == 19:  # Table
                                table_count += 1
                
                metadata["image_count"] = image_count
                metadata["text_box_count"] = text_box_count
                metadata["chart_count"] = chart_count
                metadata["table_count"] = table_count
                
            except Exception as e:
                # If metadata extraction fails, continue with basic metadata
                self.logger.warning(f"Failed to extract detailed metadata: {str(e)}")  # Use instance logger
                metadata["metadata_extraction_error"] = str(e)
            
            converter = MarkItDown()  # Instantiate the basic converter directly
            
            # Convert the presentation to markdown
            self.logger.info("Converting PPTX to markdown")  # Use instance logger
            result = converter.convert(self.file_path)
            markdown_content = result.text_content

            # Integrate image captioning
            if self.llm:
                self.logger.info("Integrating image captions into markdown")
                from pptx import Presentation  # Import Presentation here
                prs = Presentation(self.file_path)  # Re-open the presentation for shape access
                for i, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # Picture
                            # Get image data
                            image_data = shape.image.blob
                            image_stream = io.BytesIO(image_data)
                            # Create a dummy stream info (MarkItDown doesn't seem to use it fully)
                            class DummyStreamInfo:
                                def __init__(self):
                                    self.mimetype = shape.image.content_type
                                    self.extension = ".jpg"  # Or guess from mimetype
                                    self.name = shape.name
                            stream_info = DummyStreamInfo()
                            
                            # Check if the image format is supported BEFORE captioning
                            from .utils import get_image_format
                            stream_info.mimetype, stream_info.extension = get_image_format(image_data)
                            if stream_info.mimetype not in ["image/png", "image/jpeg", "image/gif", "image/webp"]:
                                self.logger.warning(f"Skipping captioning for unsupported image format: {stream_info.mimetype} (Shape: {shape.name})")
                                continue # Skip to the next shape

                            self.logger.info(f"Captioning image: {shape.name} (Format: {stream_info.mimetype})")  # Updated logging
                            # Generate caption
                            caption = langchain_caption_adapter(
                                file_stream=image_stream,  # Pass the image data directly
                                stream_info=stream_info,
                                client=self.llm,
                                model=None,
                                prompt=self.prompt
                            )
                            if caption:
                                self.logger.info(f"Generated caption: {caption[:50]}...")  # Log caption (first 50 chars)  # Use instance logger
                                # Find the image alt text placeholder in the markdown and replace it
                                # This is a HACK, as we're making assumptions about MarkItDown's output format.  It may break if MarkItDown changes.
                                # A more robust solution would require modifying MarkItDown itself.
                                alt_text_placeholder = f"![]({shape.name}.jpg)"  # Assuming MarkItDown uses shape.name.jpg as the placeholder
                                markdown_content = markdown_content.replace(alt_text_placeholder, f"![{caption}]()")
                            else:
                                self.logger.info("No caption generated")  # Log if no caption  # Use instance logger
            
            self.logger.info(f"Conversion complete, markdown content length: {len(markdown_content)} characters")

            # Define default headers to split on if not provided
            if headers_to_split_on is None:
                headers_to_split_on = [
                    ("#", "Header 1"),
                    ("##", "Header 2"),
                    ("###", "Header 3"),
                ]

            if self.split_by_page:
                # Split by slide number indicators
                documents = []
                slide_pattern = r"\n\n<!-- Slide number: (\d+) -->\n"
                slide_splits = re.split(slide_pattern, markdown_content)  # Use the modified markdown_content
                
                # The first element will be the content before the first slide indicator
                current_page_content = slide_splits[0]
                current_page_num = 1

                for i in range(1, len(slide_splits), 2):
                    if current_page_content.strip():
                        page_metadata = metadata.copy()
                        page_metadata["page_number"] = current_page_num
                        page_metadata["content_type"] = "presentation_slide"
                        
                        markdown_splitter = MarkdownHeaderTextSplitter(
                            headers_to_split_on=headers_to_split_on,
                            return_each_line=True
                        )
                        page_splits = markdown_splitter.split_text(current_page_content)
                        
                        for split in page_splits:
                            split.metadata.update(page_metadata)
                            documents.append(split)

                    current_page_num = int(slide_splits[i])
                    current_page_content = slide_splits[i + 1]

                if current_page_content.strip():
                    page_metadata = metadata.copy()
                    page_metadata["page_number"] = current_page_num
                    page_metadata["content_type"] = "presentation_slide"
                    markdown_splitter = MarkdownHeaderTextSplitter(
                            headers_to_split_on=headers_to_split_on,
                            return_each_line=True
                        )
                    page_splits = markdown_splitter.split_text(current_page_content)
                    for split in page_splits:
                        split.metadata.update(page_metadata)
                        documents.append(split)
            else:
                metadata["content_type"] = "presentation_full"
                return [Document(page_content=markdown_content, metadata=metadata)]  # Use the modified markdown_content

            return documents

        except Exception as e:
            self.logger.error(f"Failed to load and convert PPTX file: {str(e)}")
            raise ValueError(f"Failed to load and convert PPTX file: {e}")
