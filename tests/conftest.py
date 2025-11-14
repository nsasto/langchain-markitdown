import os
from typing import Any, Callable, Dict, List, Optional

import pytest
from PIL import Image


class FakeConversionResult:
    """Lightweight stand-in for MarkItDown's conversion response."""

    def __init__(
        self,
        *,
        text_content: str = "",
        pages: Optional[List[str]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        attachments: Optional[List[Any]] = None,
        output_type: str = "markdown",
        format_type: str = "text/plain",
    ):
        self.text_content = text_content
        self.pages = pages
        self.metadata = metadata or {}
        self.attachments = attachments or []
        self.output_type = output_type
        self.format_type = format_type


class FakeMarkItDown:
    """Simple callable converter used to isolate tests from the real dependency."""

    def __init__(self, result_factory: Callable[[str], FakeConversionResult]):
        self._result_factory = result_factory
        self.convert_calls: List[str] = []

    def convert(self, file_path: str) -> FakeConversionResult:
        self.convert_calls.append(file_path)
        return self._result_factory(file_path)


@pytest.fixture
def fake_converter_factory():
    """Factory fixture to generate fake converters for tests."""

    def _factory(
        *,
        text_content: str = "",
        pages: Optional[List[str]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        attachments: Optional[List[Any]] = None,
        output_type: str = "markdown",
        format_type: str = "text/plain",
        result_factory: Optional[Callable[[str], FakeConversionResult]] = None,
        exception: Optional[Exception] = None,
    ) -> FakeMarkItDown:
        if exception is not None:
            def _raise(_: str) -> FakeConversionResult:
                raise exception

            return FakeMarkItDown(_raise)

        if result_factory is not None:
            return FakeMarkItDown(result_factory)

        result = FakeConversionResult(
            text_content=text_content,
            pages=pages,
            metadata=metadata,
            attachments=attachments,
            output_type=output_type,
            format_type=format_type,
        )
        return FakeMarkItDown(lambda _: result)

    return _factory


@pytest.fixture(scope="module")
def test_docx_file(tmpdir_factory):
    """Creates a temporary DOCX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.docx")
    # Create a minimal DOCX file
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument()
        doc.add_paragraph("This is a test document.")
        doc.save(fn)
    except ImportError:
        pytest.skip("docx package not installed. Install with 'pip install python-docx'")
    return str(fn)

@pytest.fixture(scope="module")
def test_pptx_file(tmpdir_factory):
    """Creates a temporary PPTX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.pptx")
    # Create a minimal PPTX file
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = "Test Presentation"
        else:
            text_box = slide.shapes.add_textbox(0, 0, 400, 100)
            text_frame = text_box.text_frame
            text_frame.text = "Test Presentation"
        prs.save(fn)
    except ImportError:
        pytest.skip("pptx package not installed. Install with 'pip install python-pptx'")
    return str(fn)

@pytest.fixture(scope="module")
def test_xlsx_file(tmpdir_factory):
    """Creates a temporary XLSX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.xlsx")
    # Create a minimal XLSX file
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = "Test Data"
        wb.save(fn)
    except ImportError:
        pytest.skip("openpyxl package not installed. Install with 'pip install openpyxl'")
    return str(fn)


@pytest.fixture(scope="module")
def test_text_file(tmpdir_factory):
    """Creates a temporary text file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.txt")
    with open(fn, 'w', encoding='utf-8') as f:
        f.write("This is a test file.\n\nIt has multiple lines.\n\n# Header 1\n\nSome content.\n\n## Header 2\n\nMore content.")
    return str(fn)


@pytest.fixture(scope="module")
def test_image_file(tmpdir_factory):
    """Creates a temporary image file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.png")
    
    # Create a simple test image
    img = Image.new('RGB', (100, 100), color = (73, 109, 137))
    img.save(fn)
    return str(fn)
